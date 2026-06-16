# -*- coding: utf-8 -*-
"""
Generateur de presentation PowerPoint - HydroPlan SIG
Schemas visuels avec minimum de texte
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

# ===== PALETTE (ton de la plateforme HydroPlan) =====
C_DARK     = RGBColor(0x1A, 0x1A, 0x2E)
C_DARK2    = RGBColor(0x2A, 0x2A, 0x44)
C_ACCENT   = RGBColor(0xF0, 0xA5, 0x00)
C_ACCENT_D = RGBColor(0xD4, 0x92, 0x0A)
C_ACCENT_S = RGBColor(0xFF, 0xD9, 0x89)
C_BG       = RGBColor(0xFE, 0xEF, 0xDE)
C_BG_SOFT  = RGBColor(0xFF, 0xF8, 0xEE)
C_LINE     = RGBColor(0xED, 0xDF, 0xC8)
C_TEXT     = RGBColor(0x1A, 0x1A, 0x2E)
C_MUTED    = RGBColor(0x6B, 0x6B, 0x80)
C_INFO     = RGBColor(0x34, 0x98, 0xDB)
C_SUCCESS  = RGBColor(0x27, 0xAE, 0x60)
C_DANGER   = RGBColor(0xE7, 0x4C, 0x3C)
C_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]


# ============================================================
# Helpers
# ============================================================
def add_bg(slide, color=C_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, SH)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.shadow.inherit = False
    return bg

def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=C_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, italic=False, font="Inter"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_box(slide, x, y, w, h, *, fill=C_WHITE, line=C_LINE, line_w=1.0,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, shadow=False):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(line_w)
    if not shadow:
        s.shadow.inherit = False
    # Set rounded radius (only relevant for rounded rectangle)
    if shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            s.adjustments[0] = 0.12
        except Exception:
            pass
    return s

def add_box_with_text(slide, x, y, w, h, text, *, fill=C_WHITE, line=C_LINE,
                      txt_color=C_TEXT, size=12, bold=False, align=PP_ALIGN.CENTER,
                      shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_w=1.0):
    s = add_box(slide, x, y, w, h, fill=fill, line=line, line_w=line_w, shape=shape)
    tf = s.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(50000)
    tf.margin_right = Emu(50000)
    tf.margin_top = Emu(30000)
    tf.margin_bottom = Emu(30000)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Inter"
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = txt_color
    return s

def add_arrow(slide, x1, y1, x2, y2, *, color=C_DARK, weight=2.0,
              kind=MSO_CONNECTOR.STRAIGHT, arrow_end=True):
    line = slide.shapes.add_connector(kind, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    if arrow_end:
        lineEl = line.line._get_or_add_ln()
        tail = etree.SubElement(lineEl, qn('a:tailEnd'))
        tail.set('type', 'triangle')
        tail.set('w', 'med')
        tail.set('h', 'med')
    return line

def add_dashed_line(slide, x1, y1, x2, y2, color=C_MUTED, weight=1.25):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    lineEl = line.line._get_or_add_ln()
    prstDash = etree.SubElement(lineEl, qn('a:prstDash'))
    prstDash.set('val', 'dash')
    return line

def add_header(slide, title, subtitle=None, num=None):
    # bandeau sombre haut
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SW, Inches(0.95))
    band.line.fill.background()
    band.fill.solid(); band.fill.fore_color.rgb = C_DARK
    band.shadow.inherit = False
    # liseré accent
    acc = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(0.95), SW, Inches(0.06))
    acc.line.fill.background()
    acc.fill.solid(); acc.fill.fore_color.rgb = C_ACCENT
    acc.shadow.inherit = False

    if num is not None:
        chip = add_box(slide, Inches(0.45), Inches(0.27), Inches(0.55), Inches(0.4),
                       fill=C_ACCENT, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        try: chip.adjustments[0] = 0.5
        except Exception: pass
        tf = chip.text_frame
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = str(num)
        r.font.bold = True; r.font.name = "Inter"; r.font.size = Pt(16); r.font.color.rgb = C_DARK

    add_text(slide, Inches(1.15), Inches(0.18), Inches(11), Inches(0.6),
             title, size=24, bold=True, color=C_WHITE, font="Inter")
    if subtitle:
        add_text(slide, Inches(1.15), Inches(0.62), Inches(11), Inches(0.35),
                 subtitle, size=12, color=C_ACCENT_S, italic=True, font="Inter")
    # footer
    fbar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.18), SW, Inches(0.32))
    fbar.line.fill.background()
    fbar.fill.solid(); fbar.fill.fore_color.rgb = C_DARK
    fbar.shadow.inherit = False
    add_text(slide, Inches(0.3), Inches(7.18), Inches(8), Inches(0.32),
             "HydroPlan SIG  -  ORMVA Tafilalet", size=10, color=C_ACCENT_S, font="Inter")
    add_text(slide, Inches(11.5), Inches(7.18), Inches(1.7), Inches(0.32),
             "Présentation PFE", size=10, color=C_ACCENT_S, font="Inter",
             align=PP_ALIGN.RIGHT)

def small_icon_circle(slide, cx, cy, d, *, fill=C_ACCENT, line=None):
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - d/2, cy - d/2, d, d)
    c.fill.solid(); c.fill.fore_color.rgb = fill
    if line is None:
        c.line.fill.background()
    else:
        c.line.color.rgb = line
    c.shadow.inherit = False
    return c


# ============================================================
# SLIDE 1 - TITRE
# ============================================================
def slide_title():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, C_DARK)

    # bandeau decoratif accent
    add_box(s, 0, Inches(3.2), SW, Inches(1.1), fill=C_ACCENT, line=None, shape=MSO_SHAPE.RECTANGLE)

    # logo simulé
    logo_x, logo_y = Inches(5.4), Inches(0.9)
    add_box(s, logo_x, logo_y, Inches(2.5), Inches(2.0), fill=C_DARK2, line=C_ACCENT, line_w=2,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    # goutte d'eau + carte
    drop = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.3), Inches(1.25), Inches(0.7), Inches(0.9))
    drop.fill.solid(); drop.fill.fore_color.rgb = C_INFO; drop.line.fill.background()
    drop.shadow.inherit = False
    # marqueur
    pin = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(6.95), Inches(1.85), Inches(0.6), Inches(0.6))
    pin.fill.solid(); pin.fill.fore_color.rgb = C_ACCENT; pin.line.fill.background()
    pin.shadow.inherit = False

    add_text(s, Inches(0.5), Inches(3.35), Inches(12.3), Inches(0.85),
             "HydroPlan SIG", size=54, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(4.45), Inches(12.3), Inches(0.5),
             "Plateforme Web SIG  -  Gestion hydraulique & agricole",
             size=20, color=C_WHITE, align=PP_ALIGN.CENTER, italic=True)

    # icones thematiques
    icons = [("DIAGNOSTIC", C_ACCENT), ("ANALYSE", C_INFO), ("CARTE", C_SUCCESS),
             ("DIMENSIONNEMENT", C_DANGER), ("PLAN D'ACTION", C_ACCENT_S)]
    bw, gap = Inches(2.1), Inches(0.18)
    total = bw*len(icons) + gap*(len(icons)-1)
    sx = (SW - total) // 2
    for i, (t, col) in enumerate(icons):
        x = sx + i*(bw + gap)
        add_box_with_text(s, x, Inches(5.4), bw, Inches(0.55), t, fill=col,
                          line=None, txt_color=C_DARK, size=11, bold=True)

    add_text(s, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.4),
             "Office Régional de Mise en Valeur Agricole du Tafilalet",
             size=14, color=C_ACCENT_S, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.35),
             "Projet de Fin d'Etudes",
             size=12, color=C_MUTED, align=PP_ALIGN.CENTER, italic=True)


# ============================================================
# SLIDE 2 - Principe general
# ============================================================
def slide_principe():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Principe général de la plateforme",
               "Données terrain  →  Traitement SIG  →  Décision", num=1)

    # 5 etapes : Terrain -> Donnees -> Analyse -> Visualisation -> Decision
    steps = [
        ("Terrain", "Ouvrages\nhydrauliques", C_INFO,    MSO_SHAPE.ROUNDED_RECTANGLE),
        ("Saisie",  "Diagnostic\n& mesures",   C_ACCENT,  MSO_SHAPE.ROUNDED_RECTANGLE),
        ("Calcul",  "Analyse SIG\n& métier",   C_DARK2,   MSO_SHAPE.ROUNDED_RECTANGLE),
        ("Carte",   "Visualisation\ndynamique", C_SUCCESS, MSO_SHAPE.ROUNDED_RECTANGLE),
        ("Décision","Plan\nd'action",          C_DANGER,  MSO_SHAPE.ROUNDED_RECTANGLE),
    ]
    n = len(steps)
    bw, bh = Inches(2.0), Inches(1.5)
    gap = Inches(0.45)
    total = bw*n + gap*(n-1)
    sx = (SW - total) // 2
    y  = Inches(2.4)

    for i, (lab, body, col, shp) in enumerate(steps):
        x = sx + i*(bw + gap)
        # icone
        ic = add_box(s, x + bw/2 - Inches(0.4), y - Inches(0.55), Inches(0.8), Inches(0.8),
                     fill=col, line=None, shape=MSO_SHAPE.OVAL)
        add_text(s, x, y - Inches(0.55), bw, Inches(0.8), str(i+1),
                 size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # boite
        bx = add_box(s, x, y, bw, bh, fill=C_WHITE, line=col, line_w=2.0)
        tf = bx.text_frame; tf.word_wrap = True
        tf.margin_top = Emu(60000); tf.margin_bottom = Emu(60000)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = lab
        r.font.bold = True; r.font.size = Pt(14); r.font.color.rgb = col; r.font.name = "Inter"
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(11); r2.font.color.rgb = C_TEXT; r2.font.name = "Inter"

        # fleche
        if i < n - 1:
            ax = x + bw
            add_arrow(s, ax, y + bh/2, ax + gap, y + bh/2, color=C_DARK, weight=2.5)

    # bloc bas - boucle continue
    add_dashed_line(s, sx + bw/2, Inches(4.6), sx + bw/2 + total - bw, Inches(4.6),
                    color=C_ACCENT, weight=1.5)
    add_text(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.5),
             "Cycle continu : suivi & mise à jour permanents",
             size=13, italic=True, color=C_ACCENT_D, align=PP_ALIGN.CENTER)

    # 3 piliers en bas
    pillars = [("Centraliser", C_ACCENT), ("Cartographier", C_INFO), ("Décider", C_SUCCESS)]
    pw = Inches(3.6); pgap = Inches(0.35)
    ptot = pw*3 + pgap*2
    psx = (SW - ptot)//2
    for i, (t, col) in enumerate(pillars):
        x = psx + i*(pw + pgap)
        add_box_with_text(s, x, Inches(5.5), pw, Inches(1.0), t,
                          fill=col, line=None, txt_color=C_WHITE, size=18, bold=True)


# ============================================================
# SLIDE 3 - Stack technique
# ============================================================
def slide_stack():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Stack technique",
               "Django  -  GeoDjango  -  PostgreSQL / PostGIS", num=2)

    # Architecture en couches
    layers = [
        ("PRESENTATION", "HTML5 / CSS3 / JS  -  Leaflet.js  -  Chart.js", C_DARK,    Inches(1.5)),
        ("APPLICATION",  "Django  -  Vues / Formulaires / Templates",       C_ACCENT,  Inches(2.6)),
        ("METIER (GIS)", "GeoDjango  -  GDAL / GEOS  -  Calculs hydrauliques", C_INFO, Inches(3.7)),
        ("DONNEES",      "PostgreSQL  +  PostGIS (géométries)",            C_SUCCESS, Inches(4.8)),
    ]
    lw = Inches(9.5); lh = Inches(0.95)
    lx = Inches(0.5)
    for i, (name, desc, col, y) in enumerate(layers):
        # boite couche
        b = add_box(s, lx, y, lw, lh, fill=C_WHITE, line=col, line_w=2.5)
        add_text(s, lx + Inches(0.25), y + Inches(0.1), Inches(4), Inches(0.4),
                 name, size=14, bold=True, color=col)
        add_text(s, lx + Inches(0.25), y + Inches(0.45), Inches(8.7), Inches(0.45),
                 desc, size=12, color=C_TEXT)
        # numero
        n = add_box(s, lx - Inches(0.55), y + Inches(0.27), Inches(0.4), Inches(0.4),
                    fill=col, line=None, shape=MSO_SHAPE.OVAL)
        add_text(s, lx - Inches(0.55), y + Inches(0.27), Inches(0.4), Inches(0.4),
                 str(i+1), size=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # fleches verticales entre couches
    for i in range(3):
        y1 = layers[i][3] + lh
        y2 = layers[i+1][3]
        add_arrow(s, lx + lw/2, y1, lx + lw/2, y2, color=C_ACCENT_D, weight=2.0)

    # encadre cote droit - logos / specifs
    rx = Inches(10.4); rw = Inches(2.55)
    ry = Inches(1.5); rh = Inches(4.25)
    add_box(s, rx, ry, rw, rh, fill=C_DARK, line=C_ACCENT, line_w=2.0)
    add_text(s, rx, ry + Inches(0.1), rw, Inches(0.4),
             "TECHNOS", size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    techs = [
        ("Python 3.x",  C_ACCENT),
        ("Django 4+",   C_ACCENT_S),
        ("GeoDjango",   C_INFO),
        ("PostgreSQL",  C_SUCCESS),
        ("PostGIS",     C_SUCCESS),
        ("Leaflet.js",  C_ACCENT_S),
        ("GDAL / GEOS", C_INFO),
    ]
    ty = ry + Inches(0.6)
    for name, col in techs:
        chip = add_box(s, rx + Inches(0.25), ty, rw - Inches(0.5), Inches(0.4),
                       fill=col, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, rx + Inches(0.25), ty, rw - Inches(0.5), Inches(0.4),
                 name, size=11, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        ty += Inches(0.5)

    # ligne flux donnees en bas
    add_text(s, Inches(0.5), Inches(6.05), Inches(12.3), Inches(0.3),
             "Flux des données",
             size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    # mini diagramme client-server-db
    y0 = Inches(6.45)
    d = Inches(0.5)
    items = [("Navigateur", C_DARK),
             ("HTTP/AJAX", None),
             ("Django", C_ACCENT),
             ("ORM/GeoDjango", None),
             ("PostGIS", C_SUCCESS)]
    bw2 = Inches(1.7); g = Inches(0.18)
    total = bw2*3 + g*4 + Inches(1.6)*2
    sx = (SW - total)//2
    # navigateur
    add_box_with_text(s, sx, y0, bw2, d, "Navigateur",
                      fill=C_DARK, line=None, txt_color=C_WHITE, size=11, bold=True)
    add_text(s, sx + bw2 + Inches(0.05), y0, Inches(1.5), d,
             "HTTP / AJAX", size=10, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_box_with_text(s, sx + bw2 + Inches(1.6), y0, bw2, d, "Django",
                      fill=C_ACCENT, line=None, txt_color=C_DARK, size=11, bold=True)
    add_text(s, sx + 2*bw2 + Inches(1.65), y0, Inches(1.5), d,
             "ORM spatial", size=10, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)
    add_box_with_text(s, sx + 2*bw2 + Inches(3.2), y0, bw2, d, "PostGIS",
                      fill=C_SUCCESS, line=None, txt_color=C_WHITE, size=11, bold=True)


# ============================================================
# SLIDE 4 - Types d'utilisateurs
# ============================================================
def slide_users():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Types d'utilisateurs",
               "3 rôles  -  3 niveaux d'accès", num=3)

    roles = [
        ("VISITEUR",  "Consultation\nseule",       C_INFO,    "fa-eye"),
        ("ÉDITEUR",   "Diagnostic\n+ Analyse",      C_ACCENT,  "fa-edit"),
        ("OPÉRATEUR", "Tous les droits\n+ Décision", C_DANGER,  "fa-user-shield"),
    ]
    bw, bh = Inches(3.4), Inches(2.4)
    gap = Inches(0.45)
    total = bw*3 + gap*2
    sx = (SW - total)//2
    y = Inches(2.0)

    # progression hauteur (montre la hierarchie)
    heights = [Inches(2.2), Inches(2.6), Inches(3.0)]
    yshift  = [Inches(0.4), Inches(0.2), Inches(0.0)]

    for i, (name, perms, col, ic) in enumerate(roles):
        x = sx + i*(bw + gap)
        yi = y + yshift[i]
        hi = heights[i]
        # carte
        card = add_box(s, x, yi, bw, hi, fill=C_WHITE, line=col, line_w=3.0)
        # bandeau top
        head = add_box(s, x, yi, bw, Inches(0.85), fill=col, line=None,
                       shape=MSO_SHAPE.RECTANGLE)
        # avatar circle
        av = add_box(s, x + bw/2 - Inches(0.4), yi - Inches(0.4), Inches(0.8), Inches(0.8),
                     fill=C_DARK, line=col, line_w=3.0, shape=MSO_SHAPE.OVAL)
        add_text(s, x, yi - Inches(0.4), bw, Inches(0.8),
                 str(i+1), size=22, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        # nom
        add_text(s, x, yi + Inches(0.05), bw, Inches(0.5),
                 name, size=18, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # permissions
        add_text(s, x, yi + Inches(1.0), bw, Inches(0.8),
                 perms, size=14, color=C_TEXT, align=PP_ALIGN.CENTER, bold=True)
        # liste icones acces
        y_acc = yi + hi - Inches(0.85)
        accs = {
            0: [("Carte", True), ("Diagnostic", False), ("Analyse", False),
                ("Dimens.", False), ("Plan d'action", False)],
            1: [("Carte", True), ("Diagnostic", True), ("Analyse", True),
                ("Dimens.", False), ("Plan d'action", False)],
            2: [("Carte", True), ("Diagnostic", True), ("Analyse", True),
                ("Dimens.", True),  ("Plan d'action", True)],
        }[i]
        # mini tableau check
        cw = bw / 5
        for j, (lab, ok) in enumerate(accs):
            cx = x + j*cw
            mark = "✓" if ok else "—"
            mark_col = C_SUCCESS if ok else C_MUTED
            add_text(s, cx, y_acc, cw, Inches(0.35),
                     mark, size=14, bold=True, color=mark_col, align=PP_ALIGN.CENTER)
            add_text(s, cx, y_acc + Inches(0.32), cw, Inches(0.35),
                     lab, size=8, color=C_MUTED, align=PP_ALIGN.CENTER)

    # legende hierarchie
    add_text(s, Inches(0.5), Inches(5.7), Inches(12.3), Inches(0.4),
             "Niveau d'accès croissant  →",
             size=14, bold=True, color=C_ACCENT_D, align=PP_ALIGN.CENTER, italic=True)

    # fleche horizontale progression
    add_arrow(s, Inches(1.5), Inches(6.1), Inches(11.8), Inches(6.1),
              color=C_ACCENT, weight=3.0)

    # ligne accueil = tout le monde
    add_box_with_text(s, Inches(4.5), Inches(6.45), Inches(4.3), Inches(0.55),
                      "Accueil & Carte publique : accessibles à tous",
                      fill=C_BG_SOFT, line=C_ACCENT_D, line_w=1.5,
                      txt_color=C_DARK, size=11)


# ============================================================
# SLIDE 5 - App Diagnostic
# ============================================================
def slide_diagnostic():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "App Diagnostic",
               "Création des périmètres  -  Suivi & évaluation des ouvrages", num=4)

    # Schema : Perimetre -> Ouvrages de tete -> Reseau -> Diagnostic
    # Etape 1 : creation perimetre
    add_box_with_text(s, Inches(0.5), Inches(1.5), Inches(2.5), Inches(0.7),
                      "PÉRIMÈTRE", fill=C_DARK, line=None,
                      txt_color=C_WHITE, size=14, bold=True)
    # carte du perimetre stylisée
    peri = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.5), Inches(2.3), Inches(2.5), Inches(1.6))
    peri.fill.solid(); peri.fill.fore_color.rgb = C_BG_SOFT
    peri.line.color.rgb = C_ACCENT; peri.line.width = Pt(2.5)
    peri.shadow.inherit = False
    # points ouvrages dans perimetre
    coords = [(1.0, 2.6), (2.2, 2.7), (1.5, 3.1), (2.5, 3.3), (0.9, 3.4), (1.9, 3.5)]
    for cx, cy in coords:
        small_icon_circle(s, Inches(cx), Inches(cy), Inches(0.18), fill=C_ACCENT)

    # fleche -> ouvrages
    add_arrow(s, Inches(3.1), Inches(3.0), Inches(4.0), Inches(3.0),
              color=C_DARK, weight=2.5)

    # Etape 2 : ouvrages categorisés
    cats = [
        ("Ouvrages\nde tête",     ["Barrage", "Seuil", "Forage", "Khettara"], C_ACCENT),
        ("Réseau\nd'irrigation",  ["Séguia", "Prise", "Mur"],                  C_INFO),
        ("Périmètre\nirrigué",    ["Parcelles", "Assolement"],                 C_SUCCESS),
    ]
    cx = Inches(4.1); cy = Inches(1.5); cw = Inches(3.0); ch = Inches(1.45)
    for i, (cat, items, col) in enumerate(cats):
        y_ = cy + i*(ch + Inches(0.15))
        # bloc categorie
        add_box(s, cx, y_, cw, ch, fill=C_WHITE, line=col, line_w=2.0)
        add_box(s, cx, y_, Inches(1.05), ch, fill=col, line=None,
                shape=MSO_SHAPE.RECTANGLE)
        add_text(s, cx, y_, Inches(1.05), ch, cat,
                 size=11, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # items
        iw = (cw - Inches(1.2)) / max(len(items), 1)
        for j, it in enumerate(items):
            ix = cx + Inches(1.15) + j*iw
            add_box_with_text(s, ix + Inches(0.05), y_ + Inches(0.35),
                              iw - Inches(0.1), ch - Inches(0.7),
                              it, fill=C_BG_SOFT, line=col, line_w=1,
                              txt_color=C_DARK, size=10, bold=True)

    # fleche vers diagnostic
    add_arrow(s, Inches(7.2), Inches(3.0), Inches(8.1), Inches(3.0),
              color=C_DARK, weight=2.5)

    # Etape 3 : diagnostic structure
    add_box(s, Inches(8.2), Inches(1.5), Inches(4.7), Inches(4.0),
            fill=C_WHITE, line=C_DANGER, line_w=2.5)
    add_box(s, Inches(8.2), Inches(1.5), Inches(4.7), Inches(0.6),
            fill=C_DANGER, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(8.2), Inches(1.5), Inches(4.7), Inches(0.6),
             "DIAGNOSTIC", size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    diag_items = [
        ("État physique",    "Bon / Moyen / Mauvais", C_SUCCESS),
        ("Fonctionnement",   "Opérationnel / Dégradé", C_ACCENT),
        ("Pathologies",      "Fissures, érosion...",  C_DANGER),
        ("Photos & relevés", "Documentation terrain",  C_INFO),
        ("Score global",     "Note de l'ouvrage",     C_DARK),
    ]
    dy = Inches(2.25)
    for lab, val, col in diag_items:
        add_box(s, Inches(8.4), dy, Inches(0.25), Inches(0.5), fill=col, line=None,
                shape=MSO_SHAPE.RECTANGLE)
        add_text(s, Inches(8.75), dy, Inches(2.0), Inches(0.5),
                 lab, size=11, bold=True, color=C_DARK)
        add_text(s, Inches(10.5), dy, Inches(2.3), Inches(0.5),
                 val, size=10, italic=True, color=C_MUTED, align=PP_ALIGN.RIGHT)
        dy += Inches(0.6)

    # Bandeau bas : workflow
    add_box(s, Inches(0.5), Inches(5.85), Inches(12.4), Inches(1.0),
            fill=C_DARK, line=C_ACCENT, line_w=2)
    add_text(s, Inches(0.5), Inches(5.85), Inches(12.4), Inches(0.35),
             "WORKFLOW",
             size=11, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    wf = ["Création\npérimètre", "Saisie\nouvrages", "Diagnostic\nterrain", "Mise à jour\ndonnées", "Suivi\n& évaluation"]
    bw2 = Inches(2.1); gap2 = Inches(0.15)
    total = bw2*len(wf) + gap2*(len(wf)-1)
    sx2 = (SW - total)//2
    for i, t in enumerate(wf):
        x = sx2 + i*(bw2 + gap2)
        add_box_with_text(s, x, Inches(6.25), bw2, Inches(0.5),
                          t, fill=C_ACCENT, line=None, txt_color=C_DARK, size=10, bold=True)
        if i < len(wf)-1:
            add_arrow(s, x + bw2, Inches(6.5), x + bw2 + gap2, Inches(6.5),
                      color=C_WHITE, weight=1.5)


# ============================================================
# SLIDE 6 - Apps metier (vue globale)
# ============================================================
def slide_apps_metier():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Apps métier",
               "Analyse hydrologique  -  Bilan ressources  -  Efficience", num=5)

    apps = [
        ("ANALYSE\nHYDROLOGIQUE",
         "Pluies / Débits\nStations  -  BV",
         C_INFO,
         [("Pluies max", "T = 10, 50, 100 ans"),
          ("Débits", "Q de pointe"),
          ("Bassins versants", "Régionalisation")]),
        ("BILAN DES\nRESSOURCES",
         "Besoins en eau\nAssolement",
         C_SUCCESS,
         [("Cultures", "Kc / Kr"),
          ("ET₀ / ETM", "Climat"),
          ("Bilan", "Offre vs Demande")]),
        ("EFFICIENCE\nRÉSEAU",
         "Pertes / Rendement\nCanaux",
         C_ACCENT,
         [("Débits", "Mesures"),
          ("Pertes", "Infiltration"),
          ("Rendement", "% du réseau")]),
    ]
    bw, bh = Inches(4.0), Inches(4.5)
    gap = Inches(0.25)
    total = bw*3 + gap*2
    sx = (SW - total)//2
    y = Inches(1.45)

    for i, (title, sub, col, items) in enumerate(apps):
        x = sx + i*(bw + gap)
        # carte
        add_box(s, x, y, bw, bh, fill=C_WHITE, line=col, line_w=2.5)
        # bandeau top
        add_box(s, x, y, bw, Inches(1.1), fill=col, line=None, shape=MSO_SHAPE.RECTANGLE)
        # icone
        add_box(s, x + Inches(0.25), y + Inches(0.2), Inches(0.7), Inches(0.7),
                fill=C_WHITE, line=None, shape=MSO_SHAPE.OVAL)
        add_text(s, x + Inches(0.25), y + Inches(0.2), Inches(0.7), Inches(0.7),
                 str(i+1), size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(1.05), y + Inches(0.15), bw - Inches(1.2), Inches(0.8),
                 title, size=15, bold=True, color=C_WHITE)

        # sous-titre
        add_text(s, x + Inches(0.2), y + Inches(1.25), bw - Inches(0.4), Inches(0.55),
                 sub, size=12, color=C_MUTED, italic=True, align=PP_ALIGN.CENTER)

        # items
        iy = y + Inches(1.95)
        for lab, val in items:
            # icone
            small_icon_circle(s, x + Inches(0.4), iy + Inches(0.27),
                              Inches(0.22), fill=col)
            add_text(s, x + Inches(0.65), iy + Inches(0.05), bw - Inches(0.9), Inches(0.3),
                     lab, size=11, bold=True, color=C_DARK)
            add_text(s, x + Inches(0.65), iy + Inches(0.3), bw - Inches(0.9), Inches(0.25),
                     val, size=9, color=C_MUTED, italic=True)
            iy += Inches(0.7)

    # Indication : ces apps consomment des données du Diagnostic
    add_text(s, Inches(0.5), Inches(6.15), Inches(12.3), Inches(0.4),
             "Toutes les apps métier s'appuient sur les données du Diagnostic",
             size=13, italic=True, color=C_ACCENT_D, align=PP_ALIGN.CENTER, bold=True)
    # Schéma simple : diagnostic -> apps -> resultats
    yL = Inches(6.65)
    add_box_with_text(s, Inches(1.5), yL, Inches(2.5), Inches(0.45),
                      "Diagnostic", fill=C_DARK, line=None,
                      txt_color=C_WHITE, size=11, bold=True)
    add_arrow(s, Inches(4.0), yL + Inches(0.22), Inches(5.0), yL + Inches(0.22),
              color=C_DARK, weight=2)
    add_box_with_text(s, Inches(5.0), yL, Inches(3.3), Inches(0.45),
                      "Apps métier (calcul)", fill=C_ACCENT, line=None,
                      txt_color=C_DARK, size=11, bold=True)
    add_arrow(s, Inches(8.3), yL + Inches(0.22), Inches(9.3), yL + Inches(0.22),
              color=C_DARK, weight=2)
    add_box_with_text(s, Inches(9.3), yL, Inches(2.5), Inches(0.45),
                      "Résultats / Carte", fill=C_SUCCESS, line=None,
                      txt_color=C_WHITE, size=11, bold=True)


# ============================================================
# SLIDE 7 - Carte dynamique (vue globale)
# ============================================================
def slide_carte():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Carte dynamique",
               "Visualisation interactive  -  Couches  -  Requêtes", num=6)

    # mock carte cote gauche
    mx, my = Inches(0.5), Inches(1.4)
    mw, mh = Inches(7.5), Inches(5.0)
    map_bg = add_box(s, mx, my, mw, mh, fill=C_BG_SOFT, line=C_DARK, line_w=2.0,
                     shape=MSO_SHAPE.RECTANGLE)
    # grille
    for i in range(1, 10):
        x_ = mx + mw * i / 10
        add_dashed_line(s, x_, my, x_, my + mh, color=C_LINE, weight=0.5)
    for i in range(1, 7):
        y_ = my + mh * i / 7
        add_dashed_line(s, mx, y_, mx + mw, y_, color=C_LINE, weight=0.5)

    # rivieres (courbes simulees par segments)
    riv_pts = [(0.6,2.2),(2.0,2.6),(3.3,3.4),(4.6,3.0),(6.0,3.8),(7.5,4.4)]
    for i in range(len(riv_pts)-1):
        x1, y1 = riv_pts[i]; x2, y2 = riv_pts[i+1]
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                    Inches(x2), Inches(y2))
        ln.line.color.rgb = C_INFO
        ln.line.width = Pt(3)

    # decoupage admin (polygones)
    add_box(s, Inches(0.8), Inches(1.7), Inches(3.3), Inches(2.0),
            fill=RGBColor(0xFF,0xF0,0xD0), line=C_ACCENT_D, line_w=1.5,
            shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.8), Inches(1.7), Inches(3.3), Inches(0.3),
             "Zone A", size=9, italic=True, color=C_ACCENT_D)
    add_box(s, Inches(4.2), Inches(1.7), Inches(3.6), Inches(2.4),
            fill=RGBColor(0xE6,0xF7,0xE6), line=C_SUCCESS, line_w=1.5,
            shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(4.2), Inches(1.7), Inches(3.6), Inches(0.3),
             "Zone B", size=9, italic=True, color=C_SUCCESS)
    add_box(s, Inches(0.8), Inches(3.85), Inches(2.5), Inches(2.4),
            fill=RGBColor(0xFD,0xE6,0xE6), line=C_DANGER, line_w=1.5,
            shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.8), Inches(3.85), Inches(2.5), Inches(0.3),
             "Zone C", size=9, italic=True, color=C_DANGER)
    add_box(s, Inches(3.45), Inches(4.25), Inches(4.4), Inches(2.0),
            fill=RGBColor(0xE6,0xF0,0xFA), line=C_INFO, line_w=1.5,
            shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(3.45), Inches(4.25), Inches(4.4), Inches(0.3),
             "Zone D", size=9, italic=True, color=C_INFO)

    # pins ouvrages
    pins = [(1.6, 2.3, C_ACCENT), (2.8, 3.0, C_ACCENT), (4.2, 2.5, C_DANGER),
            (5.5, 3.6, C_SUCCESS), (6.7, 4.0, C_ACCENT), (1.5, 4.5, C_INFO),
            (3.4, 5.3, C_SUCCESS), (5.5, 5.5, C_DANGER), (6.8, 5.6, C_ACCENT)]
    for cx, cy, col in pins:
        small_icon_circle(s, Inches(cx), Inches(cy), Inches(0.3), fill=col)

    # canaux : largeur selon debit
    # canal 1 fin
    c1 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Inches(2.0), Inches(5.0), Inches(5.5), Inches(5.2))
    c1.line.color.rgb = C_DARK; c1.line.width = Pt(2)
    # canal 2 moyen
    c2 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Inches(2.6), Inches(2.5), Inches(5.0), Inches(2.4))
    c2.line.color.rgb = C_DARK; c2.line.width = Pt(4)
    # canal 3 large
    c3 = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Inches(0.9), Inches(3.6), Inches(6.5), Inches(3.4))
    c3.line.color.rgb = C_DARK; c3.line.width = Pt(7)

    # titre carte
    add_text(s, mx, my - Inches(0.3), mw, Inches(0.3),
             "Carte interactive (Leaflet + GeoDjango)",
             size=10, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    # cote droit : panneau controle
    px, py = Inches(8.3), Inches(1.4)
    pw, ph = Inches(4.7), Inches(5.7)
    add_box(s, px, py, pw, ph, fill=C_WHITE, line=C_LINE, line_w=1.5)

    add_text(s, px + Inches(0.2), py + Inches(0.15), pw, Inches(0.4),
             "Couches affichables", size=13, bold=True, color=C_DARK)

    layers = [
        ("Ouvrages",            C_ACCENT,  True),
        ("Périmètres",          C_INFO,    True),
        ("Canaux & séguias",    C_DARK,    True),
        ("Découpage admin.",    C_SUCCESS, True),
        ("Score / état",        C_DANGER,  False),
        ("Couches d'analyse",   C_ACCENT_S, False),
    ]
    ly = py + Inches(0.65)
    for lab, col, on in layers:
        # case a cocher
        chk = add_box(s, px + Inches(0.2), ly, Inches(0.3), Inches(0.3),
                      fill=C_WHITE, line=C_DARK, line_w=1.2,
                      shape=MSO_SHAPE.RECTANGLE)
        if on:
            add_text(s, px + Inches(0.2), ly - Inches(0.02), Inches(0.3), Inches(0.3),
                     "✓", size=14, bold=True, color=C_SUCCESS, align=PP_ALIGN.CENTER)
        # carré couleur
        add_box(s, px + Inches(0.6), ly + Inches(0.05),
                Inches(0.2), Inches(0.2), fill=col, line=None,
                shape=MSO_SHAPE.RECTANGLE)
        add_text(s, px + Inches(0.85), ly, Inches(3.5), Inches(0.3),
                 lab, size=11, color=C_TEXT)
        ly += Inches(0.45)

    # legende sous-section
    add_text(s, px + Inches(0.2), ly + Inches(0.15), pw, Inches(0.3),
             "Légende  -  Canaux (épaisseur = débit)",
             size=11, bold=True, color=C_DARK)
    ly2 = ly + Inches(0.55)
    canal_legend = [("Faible débit",  2),
                    ("Moyen débit",   4),
                    ("Fort débit",    7)]
    for lab, w in canal_legend:
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    px + Inches(0.3), ly2 + Inches(0.15),
                                    px + Inches(1.5), ly2 + Inches(0.15))
        ln.line.color.rgb = C_DARK; ln.line.width = Pt(w)
        add_text(s, px + Inches(1.7), ly2 - Inches(0.03), Inches(2.5), Inches(0.3),
                 lab, size=10, color=C_TEXT)
        ly2 += Inches(0.35)


# ============================================================
# SLIDE 8 - Requetes simples / evoluees
# ============================================================
def slide_queries():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Requêtes  -  Simples & Évoluées",
               "Filtrage  -  Croisement  -  Analyse spatiale", num=7)

    # 2 colonnes
    # Colonne 1 : Simples
    add_box(s, Inches(0.5), Inches(1.5), Inches(6.1), Inches(5.5),
            fill=C_WHITE, line=C_INFO, line_w=2.5)
    add_box(s, Inches(0.5), Inches(1.5), Inches(6.1), Inches(0.7),
            fill=C_INFO, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.5), Inches(1.5), Inches(6.1), Inches(0.7),
             "REQUÊTES SIMPLES", size=16, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    # schema simple : 1 critère -> resultat
    add_text(s, Inches(0.8), Inches(2.4), Inches(5.5), Inches(0.3),
             "1 critère  →  Liste filtrée",
             size=12, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    # schema visuel
    add_box_with_text(s, Inches(0.8), Inches(2.9), Inches(1.8), Inches(0.7),
                      "Filtre", fill=C_INFO, line=None,
                      txt_color=C_WHITE, size=12, bold=True)
    add_arrow(s, Inches(2.7), Inches(3.25), Inches(3.4), Inches(3.25),
              color=C_DARK, weight=2)
    add_box_with_text(s, Inches(3.4), Inches(2.9), Inches(2.9), Inches(0.7),
                      "Résultats", fill=C_DARK, line=None,
                      txt_color=C_WHITE, size=12, bold=True)

    # exemples
    examples_simple = [
        "Ouvrages par type",
        "Ouvrages par état (bon / mauvais)",
        "Ouvrages par périmètre",
        "Ouvrages par commune / cercle",
        "Canaux par débit",
    ]
    ey = Inches(3.85)
    for ex in examples_simple:
        add_box(s, Inches(0.95), ey, Inches(0.15), Inches(0.3),
                fill=C_INFO, line=None, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, Inches(1.2), ey, Inches(5.2), Inches(0.3),
                 ex, size=11, color=C_TEXT)
        ey += Inches(0.45)

    # Colonne 2 : Evoluees
    add_box(s, Inches(6.8), Inches(1.5), Inches(6.1), Inches(5.5),
            fill=C_WHITE, line=C_DANGER, line_w=2.5)
    add_box(s, Inches(6.8), Inches(1.5), Inches(6.1), Inches(0.7),
            fill=C_DANGER, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(6.8), Inches(1.5), Inches(6.1), Inches(0.7),
             "REQUÊTES ÉVOLUÉES", size=16, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)

    add_text(s, Inches(7.1), Inches(2.4), Inches(5.5), Inches(0.3),
             "Plusieurs critères  +  Spatial  →  Croisements",
             size=12, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    # diagramme avec plusieurs filtres convergeant
    fx, fy = Inches(7.2), Inches(2.95)
    add_box_with_text(s, fx, fy, Inches(1.2), Inches(0.45), "Critère 1",
                      fill=C_DANGER, line=None, txt_color=C_WHITE, size=10, bold=True)
    add_box_with_text(s, fx, fy + Inches(0.55), Inches(1.2), Inches(0.45), "Critère 2",
                      fill=C_DANGER, line=None, txt_color=C_WHITE, size=10, bold=True)
    add_box_with_text(s, fx, fy + Inches(1.1), Inches(1.2), Inches(0.45), "Spatial",
                      fill=C_INFO, line=None, txt_color=C_WHITE, size=10, bold=True)
    # convergent
    cx_, cy_ = Inches(10.0), Inches(3.5)
    add_box(s, cx_, cy_, Inches(1.0), Inches(1.0), fill=C_ACCENT, line=None,
            shape=MSO_SHAPE.OVAL)
    add_text(s, cx_, cy_, Inches(1.0), Inches(1.0), "ET",
             size=20, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    add_arrow(s, fx + Inches(1.2), fy + Inches(0.22), cx_, cy_ + Inches(0.5),
              color=C_DARK, weight=1.5)
    add_arrow(s, fx + Inches(1.2), fy + Inches(0.77), cx_, cy_ + Inches(0.5),
              color=C_DARK, weight=1.5)
    add_arrow(s, fx + Inches(1.2), fy + Inches(1.32), cx_, cy_ + Inches(0.5),
              color=C_DARK, weight=1.5)
    # ->  resultats
    add_arrow(s, cx_ + Inches(1.0), cy_ + Inches(0.5), Inches(11.8), cy_ + Inches(0.5),
              color=C_DARK, weight=2.5)
    add_box_with_text(s, Inches(11.8), cy_ + Inches(0.2), Inches(1.0), Inches(0.6),
                      "Carte", fill=C_DARK, line=None, txt_color=C_WHITE,
                      size=10, bold=True)

    examples_adv = [
        "Ouvrages en mauvais état dans une commune",
        "Canaux à fort débit + faible efficience",
        "Ouvrages dans rayon X km d'un seuil",
        "Périmètres à déficit hydrique",
    ]
    ey = Inches(5.0)
    for ex in examples_adv:
        add_box(s, Inches(7.1), ey, Inches(0.15), Inches(0.3),
                fill=C_DANGER, line=None, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, Inches(7.35), ey, Inches(5.5), Inches(0.3),
                 ex, size=11, color=C_TEXT)
        ey += Inches(0.4)


# ============================================================
# SLIDE 9 - Score des ouvrages
# ============================================================
def slide_score():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Score des ouvrages",
               "Évaluation visuelle  -  ÉTAT uniquement (pas dimensionnement)", num=8)

    # Echelle de score 0 - 100
    sx_ = Inches(1.0); sy_ = Inches(2.0); sw_ = Inches(11.3); sh_ = Inches(1.0)
    # gradient simulé par 5 boites
    grades = [
        ("0-20",   "Critique", C_DANGER),
        ("21-40",  "Mauvais",  RGBColor(0xE6, 0x7E, 0x22)),
        ("41-60",  "Moyen",    C_ACCENT),
        ("61-80",  "Bon",      RGBColor(0x2E, 0xCC, 0x71)),
        ("81-100", "Excellent",C_SUCCESS),
    ]
    bw = sw_ // 5
    for i, (rng, lab, col) in enumerate(grades):
        x = sx_ + bw*i
        add_box(s, x, sy_, bw, sh_, fill=col, line=C_WHITE, line_w=1, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, x, sy_, bw, Inches(0.45),
                 lab, size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, sy_ + Inches(0.45), bw, Inches(0.45),
                 rng, size=12, color=C_WHITE, align=PP_ALIGN.CENTER)

    # fleche progression
    add_arrow(s, sx_, Inches(3.25), sx_ + sw_, Inches(3.25),
              color=C_DARK, weight=2.5)
    add_text(s, sx_, Inches(3.3), sw_, Inches(0.3),
             "Score croissant  →",
             size=11, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Critères d'évaluation (entrées)
    add_text(s, Inches(0.5), Inches(3.95), Inches(12.4), Inches(0.4),
             "Critères évalués (entrées)",
             size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    crits = [
        ("Structure",      "Fissures, érosion, déformation"),
        ("Fonctionnement", "Capacité de service"),
        ("Maintenance",    "Historique d'entretien"),
        ("Sécurité",       "Risques associés"),
    ]
    cw = Inches(2.8); cgap = Inches(0.2)
    ctotal = cw*4 + cgap*3
    csx = (SW - ctotal)//2
    cy = Inches(4.5)
    for i, (lab, sub) in enumerate(crits):
        x = csx + i*(cw + cgap)
        add_box(s, x, cy, cw, Inches(1.1), fill=C_WHITE, line=C_DARK, line_w=1.5)
        add_box(s, x, cy, cw, Inches(0.4), fill=C_DARK, line=None, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, x, cy, cw, Inches(0.4),
                 lab, size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        add_text(s, x + Inches(0.1), cy + Inches(0.45), cw - Inches(0.2), Inches(0.65),
                 sub, size=10, color=C_MUTED, italic=True, align=PP_ALIGN.CENTER)

    # fleches convergent vers score
    for i in range(4):
        x = csx + i*(cw + cgap) + cw/2
        add_arrow(s, x, Inches(5.65), Inches(6.65), Inches(6.1),
                  color=C_ACCENT, weight=1.5)

    # boite score final
    add_box_with_text(s, Inches(5.4), Inches(6.1), Inches(2.5), Inches(0.8),
                      "SCORE", fill=C_ACCENT, line=C_DARK, line_w=2.5,
                      txt_color=C_DARK, size=18, bold=True)

    # Note importante
    add_box_with_text(s, Inches(8.2), Inches(6.1), Inches(4.5), Inches(0.8),
                      "L'ÉTAT, pas le dimensionnement",
                      fill=C_DANGER, line=None, txt_color=C_WHITE, size=12, bold=True)


# ============================================================
# SLIDE 10 - Decoupage administratif
# ============================================================
def slide_admin():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Découpage administratif",
               "Filtrer  -  Agréger  -  Localiser", num=9)

    # Hierarchie verticale (pyramide inverse)
    levels = [
        ("RÉGION",   "Drâa-Tafilalet",        C_DARK,    Inches(12.0)),
        ("PROVINCE", "Errachidia / Tinghir",  C_ACCENT,  Inches(9.5)),
        ("CERCLE",   "Plusieurs cercles",     C_INFO,    Inches(7.0)),
        ("COMMUNE",  "Communes rurales",      C_SUCCESS, Inches(4.5)),
        ("PÉRIMÈTRE","Unité opérationnelle",   C_DANGER,  Inches(2.5)),
    ]
    y = Inches(1.4)
    for i, (name, sub, col, w) in enumerate(levels):
        x = (SW - w)//2
        h = Inches(0.75)
        add_box(s, x, y, w, h, fill=col, line=C_WHITE, line_w=1.5)
        # nom
        add_text(s, x, y, Inches(2.5), h, name,
                 size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # sub
        add_text(s, x + Inches(2.6), y, w - Inches(2.6), h,
                 sub, size=12, color=C_WHITE, align=PP_ALIGN.CENTER, italic=True)
        # niveau
        add_box(s, x - Inches(0.6), y + Inches(0.18), Inches(0.4), Inches(0.4),
                fill=C_DARK, line=col, line_w=2, shape=MSO_SHAPE.OVAL)
        add_text(s, x - Inches(0.6), y + Inches(0.18), Inches(0.4), Inches(0.4),
                 str(i+1), size=10, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        if i < len(levels)-1:
            add_arrow(s, SW//2, y + h, SW//2, y + h + Inches(0.2),
                      color=C_DARK, weight=1.5)
        y += Inches(0.95)

    # bandeau bas - utilites
    add_box(s, Inches(0.5), Inches(6.0), Inches(12.4), Inches(1.0),
            fill=C_DARK, line=C_ACCENT, line_w=2)
    add_text(s, Inches(0.5), Inches(6.05), Inches(12.4), Inches(0.35),
             "Utilité du découpage",
             size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    uses = ["Filtrer la carte", "Agréger les indicateurs", "Croiser avec ouvrages"]
    uw = Inches(3.8); ug = Inches(0.2)
    utot = uw*3 + ug*2
    usx = (SW - utot)//2
    for i, u in enumerate(uses):
        x = usx + i*(uw + ug)
        add_box_with_text(s, x, Inches(6.4), uw, Inches(0.5),
                          u, fill=C_ACCENT, line=None, txt_color=C_DARK,
                          size=11, bold=True)


# ============================================================
# SLIDE 11 - Visualisation ouvrages (epaisseur canal selon debit)
# ============================================================
def slide_visu():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Visualisation symbolisée",
               "Ex. : épaisseur du canal  ∝  débit", num=10)

    # gauche : 3 cas (debit faible/moyen/fort)
    cases = [
        ("Faible débit",  "Q < 0.5 m³/s",  2,  C_INFO),
        ("Moyen débit",   "0.5 ≤ Q < 2",   5,  C_ACCENT),
        ("Fort débit",    "Q ≥ 2 m³/s",    10, C_DANGER),
    ]
    cw = Inches(3.5)
    cgap = Inches(0.25)
    ctot = cw*3 + cgap*2
    cstart = (SW - ctot)//2
    cy = Inches(1.6)
    for i, (lab, val, weight, col) in enumerate(cases):
        x = cstart + i*(cw + cgap)
        # carte
        add_box(s, x, cy, cw, Inches(3.0), fill=C_WHITE, line=col, line_w=2.5)
        add_box(s, x, cy, cw, Inches(0.55), fill=col, line=None, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, x, cy, cw, Inches(0.55),
                 lab, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(0.6), cw, Inches(0.3),
                 val, size=11, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)
        # mini carte avec canal
        # fond
        add_box(s, x + Inches(0.3), cy + Inches(1.05),
                cw - Inches(0.6), Inches(1.7),
                fill=C_BG_SOFT, line=C_LINE, line_w=1, shape=MSO_SHAPE.RECTANGLE)
        # canal
        ln = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                    x + Inches(0.5), cy + Inches(1.4),
                                    x + cw - Inches(0.5), cy + Inches(2.5))
        ln.line.color.rgb = C_DARK; ln.line.width = Pt(weight)
        # pin source
        small_icon_circle(s, x + Inches(0.5), cy + Inches(1.4),
                          Inches(0.25), fill=col)
        # pin sortie
        small_icon_circle(s, x + cw - Inches(0.5), cy + Inches(2.5),
                          Inches(0.25), fill=col)
        # epaisseur indicateur
        add_text(s, x, cy + Inches(2.8), cw, Inches(0.3),
                 f"Épaisseur : {weight} pt",
                 size=11, bold=True, color=col, align=PP_ALIGN.CENTER)

    # Bandeau bas - autres dimensions visuelles
    add_text(s, Inches(0.5), Inches(4.9), Inches(12.4), Inches(0.4),
             "Autres dimensions visuelles utilisées",
             size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

    dims = [
        ("Couleur",   "État de l'ouvrage",  C_ACCENT),
        ("Taille",    "Importance / capacité", C_INFO),
        ("Forme",     "Type d'ouvrage",     C_SUCCESS),
        ("Étiquette", "Nom ou ID",          C_DANGER),
        ("Opacité",   "Confiance / actif",  C_DARK),
    ]
    dw = Inches(2.35); dg = Inches(0.15)
    dt = dw*5 + dg*4
    dsx = (SW - dt)//2
    dyy = Inches(5.4)
    for lab, sub, col in dims:
        x = dsx
        add_box(s, x, dyy, dw, Inches(1.4),
                fill=C_WHITE, line=col, line_w=2)
        # carre couleur
        add_box(s, x + dw/2 - Inches(0.3), dyy + Inches(0.2),
                Inches(0.6), Inches(0.6),
                fill=col, line=None, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, x, dyy + Inches(0.85), dw, Inches(0.3),
                 lab, size=12, bold=True, color=col, align=PP_ALIGN.CENTER)
        add_text(s, x, dyy + Inches(1.12), dw, Inches(0.28),
                 sub, size=9, color=C_MUTED, italic=True, align=PP_ALIGN.CENTER)
        dsx += dw + dg


# ============================================================
# SLIDE 12 - App Dimensionnement (overview)
# ============================================================
def slide_dim_overview():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "App Dimensionnement",
               "Conception  -  Calage  -  Vérification", num=11)

    # Schema : entrees -> calcul -> sorties
    # Entrees
    add_box(s, Inches(0.5), Inches(1.5), Inches(3.0), Inches(4.8),
            fill=C_WHITE, line=C_INFO, line_w=2.5)
    add_box(s, Inches(0.5), Inches(1.5), Inches(3.0), Inches(0.6),
            fill=C_INFO, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(0.5), Inches(1.5), Inches(3.0), Inches(0.6),
             "ENTRÉES", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    inputs = [
        ("Débits", "Q de projet"),
        ("Hauteur d'eau", "Plan d'eau"),
        ("Topographie", "Profil terrain"),
        ("Sol / fondation", "Caractéristiques"),
        ("Contraintes", "Hydrauliques"),
    ]
    iy = Inches(2.25)
    for lab, sub in inputs:
        small_icon_circle(s, Inches(0.8), iy + Inches(0.3),
                          Inches(0.22), fill=C_INFO)
        add_text(s, Inches(1.05), iy, Inches(2.5), Inches(0.3),
                 lab, size=11, bold=True, color=C_DARK)
        add_text(s, Inches(1.05), iy + Inches(0.3), Inches(2.5), Inches(0.25),
                 sub, size=9, italic=True, color=C_MUTED)
        iy += Inches(0.7)

    # Fleche
    add_arrow(s, Inches(3.6), Inches(3.9), Inches(4.45), Inches(3.9),
              color=C_DARK, weight=3)

    # Module de calcul
    add_box(s, Inches(4.5), Inches(2.0), Inches(4.3), Inches(3.8),
            fill=C_DARK, line=C_ACCENT, line_w=2.5)
    add_text(s, Inches(4.5), Inches(2.15), Inches(4.3), Inches(0.45),
             "MODULE DE CALCUL", size=14, bold=True, color=C_ACCENT,
             align=PP_ALIGN.CENTER)
    # icones engrenages
    add_box(s, Inches(5.0), Inches(2.85), Inches(3.3), Inches(2.8),
            fill=C_DARK2, line=C_ACCENT, line_w=1.5)
    modules = [
        "Calage hydraulique",
        "Géométrie / forme",
        "Fondation",
        "Stabilité",
        "Vérifications",
    ]
    my_ = Inches(3.0)
    for m in modules:
        add_box(s, Inches(5.2), my_, Inches(2.9), Inches(0.42),
                fill=C_ACCENT, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        add_text(s, Inches(5.2), my_, Inches(2.9), Inches(0.42),
                 m, size=11, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
        my_ += Inches(0.52)

    # Fleche
    add_arrow(s, Inches(8.9), Inches(3.9), Inches(9.75), Inches(3.9),
              color=C_DARK, weight=3)

    # Sorties
    add_box(s, Inches(9.85), Inches(1.5), Inches(3.0), Inches(4.8),
            fill=C_WHITE, line=C_SUCCESS, line_w=2.5)
    add_box(s, Inches(9.85), Inches(1.5), Inches(3.0), Inches(0.6),
            fill=C_SUCCESS, line=None, shape=MSO_SHAPE.RECTANGLE)
    add_text(s, Inches(9.85), Inches(1.5), Inches(3.0), Inches(0.6),
             "SORTIES", size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    outputs = [
        ("Dimensions", "Largeur, hauteur"),
        ("Plan / coupe", "Schémas"),
        ("Stabilité", "OK / NON OK"),
        ("Note de calcul", "Rapport"),
        ("Estimation coût", "Quantitatif"),
    ]
    oy = Inches(2.25)
    for lab, sub in outputs:
        small_icon_circle(s, Inches(10.15), oy + Inches(0.3),
                          Inches(0.22), fill=C_SUCCESS)
        add_text(s, Inches(10.4), oy, Inches(2.5), Inches(0.3),
                 lab, size=11, bold=True, color=C_DARK)
        add_text(s, Inches(10.4), oy + Inches(0.3), Inches(2.5), Inches(0.25),
                 sub, size=9, italic=True, color=C_MUTED)
        oy += Inches(0.7)

    # Bandeau bas - rôle opérateur
    add_box_with_text(s, Inches(2.5), Inches(6.5), Inches(8.3), Inches(0.5),
                      "Accessible à l'OPÉRATEUR uniquement",
                      fill=C_DANGER, line=None, txt_color=C_WHITE,
                      size=12, bold=True)


# ============================================================
# SLIDE 13 - Dimensionnement Seuils (détail)
# ============================================================
def slide_dim_seuils():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Dimensionnement des seuils",
               "Calage  -  Forme  -  Fondation  -  Stabilité", num=12)

    # Schema coupe transversale d'un seuil + 4 etapes autour
    # zone centrale : coupe technique
    cx, cy = Inches(4.0), Inches(2.3)
    cw, ch = Inches(5.3), Inches(3.4)
    add_box(s, cx, cy, cw, ch, fill=C_BG_SOFT, line=C_DARK, line_w=1.5,
            shape=MSO_SHAPE.RECTANGLE)

    # eau amont (gauche)
    add_box(s, cx + Inches(0.2), cy + Inches(0.6), Inches(1.8), Inches(1.3),
            fill=RGBColor(0xAE, 0xD6, 0xF1), line=None, shape=MSO_SHAPE.RECTANGLE)
    # surface eau
    surf = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                  cx + Inches(0.2), cy + Inches(0.6),
                                  cx + Inches(2.0), cy + Inches(0.6))
    surf.line.color.rgb = C_INFO; surf.line.width = Pt(2)

    # seuil (corps)
    add_box(s, cx + Inches(2.0), cy + Inches(0.4), Inches(1.0), Inches(1.9),
            fill=C_ACCENT, line=C_DARK, line_w=1.5, shape=MSO_SHAPE.RECTANGLE)

    # crete
    add_text(s, cx + Inches(2.0), cy + Inches(0.05), Inches(1.0), Inches(0.35),
             "Crête", size=9, italic=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # nappe aval (descendant)
    nappe = s.shapes.add_shape(MSO_SHAPE.RIGHT_TRIANGLE,
                               cx + Inches(3.0), cy + Inches(0.4),
                               Inches(1.5), Inches(1.0))
    nappe.fill.solid(); nappe.fill.fore_color.rgb = RGBColor(0xAE, 0xD6, 0xF1)
    nappe.line.fill.background(); nappe.shadow.inherit = False

    # bassin de dissipation
    add_box(s, cx + Inches(3.0), cy + Inches(1.7), Inches(2.0), Inches(0.6),
            fill=RGBColor(0xAE, 0xD6, 0xF1), line=C_INFO, line_w=1.5,
            shape=MSO_SHAPE.RECTANGLE)
    add_text(s, cx + Inches(3.0), cy + Inches(2.32), Inches(2.0), Inches(0.25),
             "Bassin de dissipation",
             size=8, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    # fondation
    add_box(s, cx + Inches(0.1), cy + Inches(2.3), Inches(5.0), Inches(0.85),
            fill=RGBColor(0xC9, 0xB3, 0x8A), line=C_DARK, line_w=1.5,
            shape=MSO_SHAPE.RECTANGLE)
    add_text(s, cx + Inches(0.1), cy + Inches(2.7), Inches(5.0), Inches(0.3),
             "Fondation / Sol",
             size=10, italic=True, color=C_DARK, align=PP_ALIGN.CENTER)

    # labels eau
    add_text(s, cx + Inches(0.3), cy + Inches(1.0), Inches(1.5), Inches(0.3),
             "Eau amont", size=9, italic=True, color=C_INFO)

    add_text(s, cx, cy - Inches(0.1), cw, Inches(0.3),
             "Coupe transversale - seuil",
             size=10, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    # 4 etiquettes autour - les etapes
    steps = [
        ("CALAGE", "Hauteur / cote", C_INFO,
         Inches(0.5), Inches(1.5), True),         # haut-gauche
        ("FORME", "Géométrie / profil", C_ACCENT,
         Inches(9.5), Inches(1.5), True),          # haut-droite
        ("FONDATION", "Sol / encrage", C_SUCCESS,
         Inches(0.5), Inches(4.3), True),         # bas-gauche
        ("STABILITÉ", "Glissement / renversement", C_DANGER,
         Inches(9.5), Inches(4.3), True),          # bas-droite
    ]
    for name, sub, col, x, y, _ in steps:
        add_box(s, x, y, Inches(3.3), Inches(1.5), fill=C_WHITE, line=col, line_w=2.5)
        add_box(s, x, y, Inches(3.3), Inches(0.55), fill=col, line=None,
                shape=MSO_SHAPE.RECTANGLE)
        add_text(s, x, y, Inches(3.3), Inches(0.55),
                 name, size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, y + Inches(0.6), Inches(3.3), Inches(0.9),
                 sub, size=11, color=C_TEXT, align=PP_ALIGN.CENTER)

    # connexions etiquettes -> schéma central
    centerx = cx + cw/2
    centery = cy + ch/2
    # haut-gauche
    add_dashed_line(s, Inches(3.8), Inches(2.25), cx + Inches(1.2), cy + Inches(0.5),
                    color=C_INFO, weight=1.5)
    # haut-droite
    add_dashed_line(s, Inches(9.5), Inches(2.25), cx + Inches(4.5), cy + Inches(0.5),
                    color=C_ACCENT, weight=1.5)
    # bas-gauche
    add_dashed_line(s, Inches(3.8), Inches(5.0), cx + Inches(1.0), cy + Inches(2.8),
                    color=C_SUCCESS, weight=1.5)
    # bas-droite
    add_dashed_line(s, Inches(9.5), Inches(5.0), cx + Inches(4.5), cy + Inches(2.8),
                    color=C_DANGER, weight=1.5)

    # Bandeau bas - étape de vérification
    add_box(s, Inches(0.5), Inches(6.3), Inches(12.4), Inches(0.7),
            fill=C_DARK, line=C_ACCENT, line_w=2)
    add_text(s, Inches(0.5), Inches(6.3), Inches(12.4), Inches(0.7),
             "Analyse de stabilité  →  Vérifications réglementaires  →  Validation",
             size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 14 - Plan d'action - definition
# ============================================================
def slide_plan_def():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Plan d'action",
               "Du diagnostic à la décision", num=13)

    # Schema : agregation -> priorisation -> plan
    # 3 colonnes verticales puis sortie globale
    cw = Inches(3.5); cgap = Inches(0.3)
    ctot = cw*3 + cgap*2
    csx = (SW - ctot)//2
    cy = Inches(1.5)
    cols = [
        ("CONSTAT", "Issues du Diagnostic",
         ["Ouvrages dégradés", "Périmètres à risque", "Pertes du réseau"], C_DANGER),
        ("PRIORISATION", "Score & criticité",
         ["Tri par criticité", "Filtrage admin.", "Estimation coût"], C_ACCENT),
        ("PLAN", "Actions à mener",
         ["Réhabilitations", "Nouvelles réalisations", "Maintenance"], C_SUCCESS),
    ]
    for i, (name, sub, items, col) in enumerate(cols):
        x = csx + i*(cw + cgap)
        add_box(s, x, cy, cw, Inches(4.3), fill=C_WHITE, line=col, line_w=2.5)
        add_box(s, x, cy, cw, Inches(0.7), fill=col, line=None, shape=MSO_SHAPE.RECTANGLE)
        add_text(s, x, cy + Inches(0.05), cw, Inches(0.4),
                 name, size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        add_text(s, x, cy + Inches(0.4), cw, Inches(0.3),
                 sub, size=10, color=C_BG_SOFT, italic=True, align=PP_ALIGN.CENTER)
        iy = cy + Inches(1.0)
        for it in items:
            add_box(s, x + Inches(0.3), iy, cw - Inches(0.6), Inches(0.7),
                    fill=C_BG_SOFT, line=col, line_w=1)
            small_icon_circle(s, x + Inches(0.55), iy + Inches(0.35),
                              Inches(0.2), fill=col)
            add_text(s, x + Inches(0.85), iy + Inches(0.18), cw - Inches(1), Inches(0.4),
                     it, size=11, bold=True, color=C_DARK)
            iy += Inches(0.85)
        # fleche vers suivant
        if i < 2:
            add_arrow(s, x + cw, cy + Inches(2.0),
                      x + cw + cgap, cy + Inches(2.0),
                      color=C_DARK, weight=2.5)

    # Bandeau bas
    add_box(s, Inches(0.5), Inches(6.2), Inches(12.4), Inches(0.85),
            fill=C_DARK, line=C_ACCENT, line_w=2)
    add_text(s, Inches(0.5), Inches(6.25), Inches(12.4), Inches(0.35),
             "Le plan d'action est piloté par l'opérateur",
             size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(6.6), Inches(12.4), Inches(0.35),
             "Visible et exportable depuis la plateforme",
             size=11, italic=True, color=C_BG_SOFT, align=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 15 - Plan d'action - outils
# ============================================================
def slide_plan_tools():
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_header(s, "Plan d'action  -  Outils",
               "Gestion  -  Suivi  -  Reporting", num=14)

    # 4 outils en 2 lignes x 2 colonnes
    tools = [
        ("Calendrier",       "Planification dans le temps",     C_INFO,    "📅"),
        ("Suivi d'avancement","Taux & jalons",                  C_ACCENT,  "📈"),
        ("Carte d'intervention","Localisation des chantiers",   C_SUCCESS, "🗺"),
        ("Reporting",        "Tableaux de bord & exports",      C_DANGER,  "📊"),
    ]
    bw, bh = Inches(5.8), Inches(2.3)
    gap = Inches(0.3)
    sx_ = (SW - (bw*2 + gap))//2
    sy_ = Inches(1.4)
    positions = [
        (sx_,           sy_),
        (sx_ + bw + gap, sy_),
        (sx_,           sy_ + bh + gap),
        (sx_ + bw + gap, sy_ + bh + gap),
    ]
    for (lab, sub, col, ic), (x, y) in zip(tools, positions):
        add_box(s, x, y, bw, bh, fill=C_WHITE, line=col, line_w=2.5)
        # bandeau coloré gauche
        add_box(s, x, y, Inches(0.4), bh, fill=col, line=None, shape=MSO_SHAPE.RECTANGLE)
        # icone (cercle avec emoji)
        add_box(s, x + Inches(0.7), y + Inches(0.5), Inches(1.3), Inches(1.3),
                fill=col, line=None, shape=MSO_SHAPE.OVAL)
        add_text(s, x + Inches(0.7), y + Inches(0.5), Inches(1.3), Inches(1.3),
                 ic, size=44, color=C_WHITE, align=PP_ALIGN.CENTER)
        # textes
        add_text(s, x + Inches(2.2), y + Inches(0.55), bw - Inches(2.4), Inches(0.7),
                 lab, size=20, bold=True, color=col)
        add_text(s, x + Inches(2.2), y + Inches(1.25), bw - Inches(2.4), Inches(0.85),
                 sub, size=13, color=C_MUTED, italic=True)

    # bandeau bas - feedback loop
    add_box(s, Inches(0.5), Inches(6.3), Inches(12.4), Inches(0.75),
            fill=C_DARK, line=C_ACCENT, line_w=2)
    add_text(s, Inches(0.5), Inches(6.35), Inches(12.4), Inches(0.3),
             "Boucle de mise à jour permanente",
             size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    # diagramme circulaire mini
    items = ["Diagnostic", "Plan", "Action", "Suivi"]
    iw = Inches(2.0); ig = Inches(0.1)
    itot = iw*4 + ig*3
    isx = (SW - itot)//2
    for i, it in enumerate(items):
        x = isx + i*(iw + ig)
        add_box_with_text(s, x, Inches(6.65), iw, Inches(0.35),
                          it, fill=C_ACCENT, line=None, txt_color=C_DARK,
                          size=10, bold=True)
        if i < 3:
            add_arrow(s, x + iw, Inches(6.82), x + iw + ig, Inches(6.82),
                      color=C_WHITE, weight=1.5)


# ============================================================
# SLIDE FINALE
# ============================================================
def slide_final():
    s = prs.slides.add_slide(BLANK)
    add_bg(s, C_DARK)

    # accent stripe
    add_box(s, 0, Inches(3.0), SW, Inches(0.1), fill=C_ACCENT, line=None,
            shape=MSO_SHAPE.RECTANGLE)

    add_text(s, Inches(0.5), Inches(2.0), Inches(12.3), Inches(0.85),
             "Merci de votre attention", size=48, bold=True, color=C_WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0.5), Inches(3.25), Inches(12.3), Inches(0.4),
             "Questions  &  Démo",
             size=20, color=C_ACCENT, align=PP_ALIGN.CENTER, italic=True)

    # mini schema en bas (recap)
    items = ["Diagnostic", "Analyse", "Carte", "Dimensionnement", "Plan"]
    bw = Inches(2.1); bg = Inches(0.15)
    btot = bw*5 + bg*4
    bsx = (SW - btot)//2
    for i, it in enumerate(items):
        x = bsx + i*(bw + bg)
        add_box_with_text(s, x, Inches(4.7), bw, Inches(0.7),
                          it, fill=C_ACCENT, line=None, txt_color=C_DARK,
                          size=13, bold=True)
        if i < 4:
            add_arrow(s, x + bw, Inches(5.05), x + bw + bg, Inches(5.05),
                      color=C_ACCENT, weight=2)

    add_text(s, Inches(0.5), Inches(6.4), Inches(12.3), Inches(0.4),
             "HydroPlan SIG  -  ORMVA Tafilalet",
             size=14, color=C_ACCENT_S, align=PP_ALIGN.CENTER, italic=True)


# ============================================================
# Build
# ============================================================
if __name__ == "__main__":
    slide_title()
    slide_principe()
    slide_stack()
    slide_users()
    slide_diagnostic()
    slide_apps_metier()
    slide_carte()
    slide_queries()
    slide_score()
    slide_admin()
    slide_visu()
    slide_dim_overview()
    slide_dim_seuils()
    slide_plan_def()
    slide_plan_tools()
    slide_final()

    out = "HydroPlan_SIG_Presentation.pptx"
    prs.save(out)
    print(f"OK -> {out}  ({len(prs.slides)} slides)")
